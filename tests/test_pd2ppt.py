import os
import pandas as pd
import pd2ppt

from pptx.util import Cm
from pd2ppt import df_to_powerpoint



def test_do_formatting_missing():
    # '' - int
    assert pd2ppt.pd2ppt._do_formatting(9, '') == '9'
    assert pd2ppt.pd2ppt._do_formatting(99, '') == '99'
    assert pd2ppt.pd2ppt._do_formatting(999, '') == '999'
    assert pd2ppt.pd2ppt._do_formatting(9999, '') == '9,999'
    assert pd2ppt.pd2ppt._do_formatting(99999, '') == '99,999'

    # '' - float
    assert pd2ppt.pd2ppt._do_formatting(.999000, '') == '0.999000'
    assert pd2ppt.pd2ppt._do_formatting(9.999000, '') == '9.999000'
    assert pd2ppt.pd2ppt._do_formatting(99.999000, '') == '99.999000'
    assert pd2ppt.pd2ppt._do_formatting(999.999000, '') == '999.999000'
    assert pd2ppt.pd2ppt._do_formatting(9999.999000, '') == '9999.999000'
    assert pd2ppt.pd2ppt._do_formatting(99999.999000, '') == '99999.999000'
    assert pd2ppt.pd2ppt._do_formatting(1.23456789, '') == '1.234568'

    # '' - string
    assert pd2ppt.pd2ppt._do_formatting('.999000', '') == '.999000'
    assert pd2ppt.pd2ppt._do_formatting('9.999000', '') == '9.999000'
    assert pd2ppt.pd2ppt._do_formatting('99.999000', '') == '99.999000'
    assert pd2ppt.pd2ppt._do_formatting('999.999000', '') == '999.999000'
    assert pd2ppt.pd2ppt._do_formatting('9999.999000', '') == '9999.999000'
    assert pd2ppt.pd2ppt._do_formatting('99999.999000', '') == '99999.999000'
    assert pd2ppt.pd2ppt._do_formatting('1.23456789', '') == '1.23456789'



def test_do_formatting_comma():
    # ',' The ',' option signals the use of a comma for a thousands separator
    assert pd2ppt.pd2ppt._do_formatting(9, ',') == '9'
    assert pd2ppt.pd2ppt._do_formatting(99, ',') == '99'
    assert pd2ppt.pd2ppt._do_formatting(999, ',') == '999'
    assert pd2ppt.pd2ppt._do_formatting(9999, ',') == '9,999'
    assert pd2ppt.pd2ppt._do_formatting(99999, ',') == '99,999'
    assert pd2ppt.pd2ppt._do_formatting(99999.999, ',') == '99,999.999'

    assert pd2ppt.pd2ppt._do_formatting(.999000, ',') == '0.999'
    assert pd2ppt.pd2ppt._do_formatting(9.999000, ',') == '9.999'
    assert pd2ppt.pd2ppt._do_formatting(99.999000, ',') == '99.999'
    assert pd2ppt.pd2ppt._do_formatting(999.999000, ',') == '999.999'
    assert pd2ppt.pd2ppt._do_formatting(9999.999000, ',') == '9,999.999'
    assert pd2ppt.pd2ppt._do_formatting(99999.999000, ',') == '99,999.999'
    assert pd2ppt.pd2ppt._do_formatting(1.23456789, ',') == '1.23456789'

    # except cases
    assert pd2ppt.pd2ppt._do_formatting('.999000', ',') == '.999000'
    assert pd2ppt.pd2ppt._do_formatting('9.999000', ',') == '9.999000'
    assert pd2ppt.pd2ppt._do_formatting('99.999000', ',') == '99.999000'
    assert pd2ppt.pd2ppt._do_formatting('999.999000', ',') == '999.999000'
    assert pd2ppt.pd2ppt._do_formatting('9999.999000', ',') == '9999.999000'
    assert pd2ppt.pd2ppt._do_formatting('99999.999000', ',') == '99999.999000'
    assert pd2ppt.pd2ppt._do_formatting('1.23456789', ',') == '1.23456789'



def test_do_formatting_decimal():
    assert pd2ppt.pd2ppt._do_formatting(9, '.') == '9'
    assert pd2ppt.pd2ppt._do_formatting(99, '.') == '99'
    assert pd2ppt.pd2ppt._do_formatting(999, '.') == '999'
    assert pd2ppt.pd2ppt._do_formatting(9999, '.') == '9999'
    assert pd2ppt.pd2ppt._do_formatting(99999, '.') == '99999'
    assert pd2ppt.pd2ppt._do_formatting(99999.999, '.') == '99999.999'

    assert pd2ppt.pd2ppt._do_formatting(.999000, '.') == '0.999'
    assert pd2ppt.pd2ppt._do_formatting(9.999000, '.') == '9.999'
    assert pd2ppt.pd2ppt._do_formatting(99.999000, '.') == '99.999'
    assert pd2ppt.pd2ppt._do_formatting(999.999000, '.') == '999.999'
    assert pd2ppt.pd2ppt._do_formatting(9999.999000, '.') == '9999.999'
    assert pd2ppt.pd2ppt._do_formatting(99999.999000, '.') == '99999.999'
    assert pd2ppt.pd2ppt._do_formatting(1.23456789, '.') == '1.23456789'

    # except cases
    assert pd2ppt.pd2ppt._do_formatting('.999000', '.') == '.999000'
    assert pd2ppt.pd2ppt._do_formatting('9.999000', '.') == '9.999000'
    assert pd2ppt.pd2ppt._do_formatting('99.999000', '.') == '99.999000'
    assert pd2ppt.pd2ppt._do_formatting('999.999000', '.') == '999.999000'
    assert pd2ppt.pd2ppt._do_formatting('9999.999000', '.') == '9999.999000'
    assert pd2ppt.pd2ppt._do_formatting('99999.999000', '.') == '99999.999000'
    assert pd2ppt.pd2ppt._do_formatting('1.23456789', '.') == '1.23456789'


def test_do_formatting_decimal_4():
    assert pd2ppt.pd2ppt._do_formatting(9, '.4') == '9'
    assert pd2ppt.pd2ppt._do_formatting(99, '.4') == '99'
    assert pd2ppt.pd2ppt._do_formatting(999, '.4') == '999'
    assert pd2ppt.pd2ppt._do_formatting(9999, '.4') == '9999'
    assert pd2ppt.pd2ppt._do_formatting(99999, '.4') == '1E+05'
    assert pd2ppt.pd2ppt._do_formatting(99999.999, '.4') == '1E+05'

    assert pd2ppt.pd2ppt._do_formatting(.999000, '.4') == '0.999'
    assert pd2ppt.pd2ppt._do_formatting(9.999000, '.4') == '9.999'
    assert pd2ppt.pd2ppt._do_formatting(99.999000, '.4') == '100'
    assert pd2ppt.pd2ppt._do_formatting(999.999000, '.4') == '1000'
    assert pd2ppt.pd2ppt._do_formatting(9999.999000, '.4') == '1E+04'
    assert pd2ppt.pd2ppt._do_formatting(99999.999000, '.4') == '1E+05'
    assert pd2ppt.pd2ppt._do_formatting(1.23456789, '.4') == '1.235'

    assert pd2ppt.pd2ppt._do_formatting('.999000', '.4') == '.999000'
    assert pd2ppt.pd2ppt._do_formatting('9.999000', '.4') == '9.999000'
    assert pd2ppt.pd2ppt._do_formatting('99.999000', '.4') == '99.999000'
    assert pd2ppt.pd2ppt._do_formatting('999.999000', '.4') == '999.999000'
    assert pd2ppt.pd2ppt._do_formatting('9999.999000', '.4') == '9999.999000'
    assert pd2ppt.pd2ppt._do_formatting('99999.999000', '.4') == '99999.999000'
    assert pd2ppt.pd2ppt._do_formatting('1.23456789', '.4') == '1.23456789'



def test_do_formatting_decimal_end_with_R():
    assert pd2ppt.pd2ppt._do_formatting(1.23456789, '.4R') == '1.23456789'
    assert pd2ppt.pd2ppt._do_formatting(12345, '.1R') == '10,000'
    assert pd2ppt.pd2ppt._do_formatting(12345, '.2R') == '12,000'
    assert pd2ppt.pd2ppt._do_formatting(12345, '.3R') == '12,300'
    assert pd2ppt.pd2ppt._do_formatting(12345, '.4R') == '12,340'



def test_process_position_parameter():
    assert pd2ppt.pd2ppt.process_position_parameter(None) == Cm(4)
    assert pd2ppt.pd2ppt.process_position_parameter(1) == Cm(1)
    assert pd2ppt.pd2ppt.process_position_parameter(100) == Cm(100)
    assert pd2ppt.pd2ppt.process_position_parameter(1.2) == 1.2
    assert pd2ppt.pd2ppt.process_position_parameter('3') == '3'



def test_integration_df_to_powerpoint():
    df = pd.DataFrame(
        {'District':['Hampshire', 'Dorset', 'Wiltshire', 'Worcestershire'],
         'Population':[25000, 500000, 735298, 12653],
         'Ratio':[1.56, 7.34, 3.67, 8.23]})

    shape = df_to_powerpoint(
                r"test1.pptx", df,
                col_formatters=['', ',', '.2'], rounding=['', 3, ''])

    assert shape.table.cell(0,0).text_frame.text == "District"
    assert shape.table.cell(0,1).text_frame.text == "Population"
    assert shape.table.cell(0,2).text_frame.text == "Ratio"

    assert shape.table.cell(1,0).text_frame.text == "Hampshire"
    assert shape.table.cell(1,1).text_frame.text == "25,000"
    assert shape.table.cell(1,2).text_frame.text == "1.6"

    assert shape.table.cell(2,0).text_frame.text == "Dorset"
    assert shape.table.cell(2,1).text_frame.text == "500,000"
    assert shape.table.cell(2,2).text_frame.text == "7.3"

    assert shape.table.cell(3,0).text_frame.text == "Wiltshire"
    assert shape.table.cell(3,1).text_frame.text == "735,298"
    assert shape.table.cell(3,2).text_frame.text == "3.7"

    assert shape.table.cell(4,0).text_frame.text == "Worcestershire"
    assert shape.table.cell(4,1).text_frame.text == "12,653"
    assert shape.table.cell(4,2).text_frame.text == "8.2"


    shape = df_to_powerpoint(
                r"test2.pptx", df,
                left=1, top=1, width=10, height=15,
                col_formatters=['', '.', '.3'], rounding=['', 3, ''])

    assert shape.table.cell(0,0).text_frame.text == "District"
    assert shape.table.cell(0,1).text_frame.text == "Population"
    assert shape.table.cell(0,2).text_frame.text == "Ratio"

    assert shape.table.cell(1,0).text_frame.text == "Hampshire"
    assert shape.table.cell(1,1).text_frame.text == "25000"
    assert shape.table.cell(1,2).text_frame.text == "1.56"

    assert shape.table.cell(2,0).text_frame.text == "Dorset"
    assert shape.table.cell(2,1).text_frame.text == "500000"
    assert shape.table.cell(2,2).text_frame.text == "7.34"

    assert shape.table.cell(3,0).text_frame.text == "Wiltshire"
    assert shape.table.cell(3,1).text_frame.text == "735298"
    assert shape.table.cell(3,2).text_frame.text == "3.67"

    assert shape.table.cell(4,0).text_frame.text == "Worcestershire"
    assert shape.table.cell(4,1).text_frame.text == "12653"
    assert shape.table.cell(4,2).text_frame.text == "8.23"
