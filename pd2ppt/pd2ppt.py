import six
import pandas as pd

from math import *
from pptx import Presentation
from pptx.util import Cm, Pt


round_to_n = lambda x, n: round(x, -int(floor(log10(abs(x)))) + (n - 1))


class TableStyle:
    """Class holding ids of table templates"""
    NoStyleNoGrid = '{2D5ABB26-0587-4C30-8999-92F81FD0307C}'
    ThemedStyle1Accent1 = '{3C2FFA5D-87B4-456A-9821-1D50468CF0F}'
    ThemedStyle1Accent2 = '{284E427A-3D55-4303-BF80-6455036E1DE7}'
    ThemedStyle1Accent3 = '{69C7853C-536D-4A76-A0AE-DD22124D55A5}'
    ThemedStyle1Accent4 = '{775DCB02-9BB8-47FD-8907-85C794F793BA}'
    ThemedStyle1Accent5 = '{35758FB7-9AC5-4552-8A53-C91805E547FA}'
    ThemedStyle1Accent6 = '{08FB837D-C827-4EFA-A057-4D05807E0F7C}'
    NoStyleTableGrid = '{5940675A-B579-460E-94D1-54222C63F5DA}'
    ThemedStyle2Accent1 = '{D113A9D2-9D6B-4929-AA2D-F23B5EE8CBE7}'
    ThemedStyle2Accent2 = '{18603FDC-E32A-4AB5-989C-0864C3EAD2B8}'
    ThemedStyle2Accent3 = '{306799F8-075E-4A3A-A7F6-7FBC6576F1A4}'
    ThemedStyle2Accent4 = '{E269D01E-BC32-4049-B463-5C60D7B0CCD2}'
    ThemedStyle2Accent5 = '{327F97BB-C833-4FB7-BDE5-3F7075034690}'
    ThemedStyle2Accent6 = '{638B1855-1B75-4FBE-930C-398BA8C253C6}'
    LightStyle1 = '{9D7B26C5-4107-4FEC-AEDC-1716B250A1EF}'
    LightStyle1Accent1 = '{3B4B98B0-60AC-42C2-AFA5-B58CD77FA1E5}'
    LightStyle1Accent2 = '{0E3FDE45-AF77-4B5C-9715-49D594BDF05E}'
    LightStyle1Accent3 = '{C083E6E3-FA7D-4D7B-A595-EF9225AFEA82}'
    LightStyle1Accent4 = '{D27102A9-8310-4765-A935-A1911B00CA55}'
    LightStyle1Accent5 = '{5FD0F851-EC5A-4D38-B0AD-8093EC10F338}'
    LightStyle1Accent6 = '{68D230F3-CF80-4859-8CE7-A43EE81993B5}'
    LightStyle2 = '{7E9639D4-E3E2-4D34-9284-5A2195B3D0D7}'
    LightStyle2Accent1 = '{69012ECD-51FC-41F1-AA8D-1B2483CD663E}'
    LightStyle2Accent2 = '{72833802-FEF1-4C79-8D5D-14CF1EAF98D9}'
    LightStyle2Accent3 = '{F2DE63D5-997A-4646-A377-4702673A728D}'
    LightStyle2Accent4 = '{17292A2E-F333-43FB-9621-5CBBE7FDCDCB}'
    LightStyle2Accent5 = '{5A111915-BE36-4E01-A7E5-04B1672EAD32}'
    LightStyle2Accent6 = '{912C8C85-51F0-491E-9774-3900AFEF0FD7}'
    LightStyle3 = '{616DA210-FB5B-4158-B5E0-FEB733F419BA}'
    LightStyle3Accent1 = '{BC89EF96-8CEA-46FF-86C4-4CE0E7609802}'
    LightStyle3Accent2 = '{5DA37D80-6434-44D0-A028-1B22A696006F}'
    LightStyle3Accent3 = '{8799B23B-EC83-4686-B30A-512413B5E67A}'
    LightStyle3Accent4 = '{ED083AE6-46FA-4A59-8FB0-9F97EB10719F}'
    LightStyle3Accent5 = '{BDBED569-4797-4DF1-A0F4-6AAB3CD982D8}'
    LightStyle3Accent6 = '{E8B1032C-EA38-4F05-BA0D-38AFFFC7BED3}'
    MediumStyle1 = '{793D81CF-94F2-401A-BA57-92F5A7B2D0C5}'
    MediumStyle1Accent1 = '{B301B821-A1FF-4177-AEE7-76D212191A09}'
    MediumStyle1Accent2 = '{9DCAF9ED-07DC-4A11-8D7F-57B35C25682E}'
    MediumStyle1Accent3 = '{1FECB4D8-DB02-4DC6-A0A2-4F2EBAE1DC90}'
    MediumStyle1Accent4 = '{1E171933-4619-4E11-9A3F-F7608DF75F80}'
    MediumStyle1Accent5 = '{FABFCF23-3B69-468F-B69F-88F6DE6A72F2}'
    MediumStyle1Accent6 = '{10A1B5D5-9B99-4C35-A422-299274C87663}'
    MediumStyle2 = '{073A0DAA-6AF3-43AB-8588-CEC1D06C72B9}'
    MediumStyle2Accent1 = '{5C22544A-7EE6-4342-B048-85BDC9FD1C3A}'
    MediumStyle2Accent2 = '{21E4AEA4-8DFA-4A89-87EB-49C32662AFE0}'
    MediumStyle2Accent3 = '{F5AB1C69-6EDB-4FF4-983F-18BD219EF322}'
    MediumStyle2Accent4 = '{00A15C55-8517-42AA-B614-E9B94910E393}'
    MediumStyle2Accent5 = '{7DF18680-E054-41AD-8BC1-D1AEF772440D}'
    MediumStyle2Accent6 = '{93296810-A885-4BE3-A3E7-6D5BEEA58F35}'
    MediumStyle3 = '{8EC20E35-A176-4012-BC5E-935CFFF8708E}'
    MediumStyle3Accent1 = '{6E25E649-3F16-4E02-A733-19D2CDBF48F0}'
    MediumStyle3Accent2 = '{85BE263C-DBD7-4A20-BB59-AAB30ACAA65A}'
    MediumStyle3Accent3 = '{EB344D84-9AFB-497E-A393-DC336BA19D2E}'
    MediumStyle3Accent4 = '{EB9631B5-78F2-41C9-869B-9F39066F8104}'
    MediumStyle3Accent5 = '{74C1A8A3-306A-4EB7-A6B1-4F7E0EB9C5D6}'
    MediumStyle3Accent6 = '{2A488322-F2BA-4B5B-9748-0D474271808F}'
    MediumStyle4 = '{D7AC3CCA-C797-4891-BE02-D94E43425B78}'
    MediumStyle4Accent1 = '{69CF1AB2-1976-4502-BF36-3FF5EA218861}'
    MediumStyle4Accent2 = '{8A107856-5554-42FB-B03E-39F5DBC370BA}'
    MediumStyle4Accent3 = '{0505E3EF-67EA-436B-97B2-0124C06EBD24}'
    MediumStyle4Accent4 = '{C4B1156A-380E-4F78-BDF5-A606A8083BF9}'
    MediumStyle4Accent5 = '{22838BEF-8BB2-4498-84A7-C5851F593DF1}'
    MediumStyle4Accent6 = '{16D9F66E-5EB9-4882-86FB-DCBF35E3C3E4}'
    DarkStyle1 = '{E8034E78-7F5D-4C2E-B375-FC64B27BC917}'
    DarkStyle1Accent1 = '{125E5076-3810-47DD-B79F-674D7AD40C01}'
    DarkStyle1Accent2 = '{37CE84F3-28C3-443E-9E96-99CF82512B78}'
    DarkStyle1Accent3 = '{D03447BB-5D67-496B-8E87-E561075AD55C}'
    DarkStyle1Accent4 = '{E929F9F4-4A8F-4326-A1B4-22849713DDAB}'
    DarkStyle1Accent5 = '{8FD4443E-F989-4FC4-A0C8-D5A2AF1F390B}'
    DarkStyle1Accent6 = '{AF606853-7671-496A-8E4F-DF71F8EC918B}'
    DarkStyle2 = '{5202B0CA-FC54-4496-8BCA-5EF66A818D29}'
    DarkStyle2Accent1Accent2 = '{0660B408-B3CF-4A94-85FC-2B1E0A45F4A2}'
    DarkStyle2Accent3Accent4 = '{91EBBBCC-DAD2-459C-BE2E-F6DE35CF9A28}'
    DarkStyle2Accent5Accent6 = '{46F890A9-2807-4EBB-B81D-B2AA78EC7F39}'


def _do_formatting(value, format_str):
    """Format value according to format_str, and deal
    sensibly with format_str if it is missing or invalid.
    """
    if format_str == '':
        if type(value) in six.integer_types:
            format_str = ','
        elif type(value) is float:
            format_str = 'f'
        elif type(value) is str:
            format_str = 's'
    elif format_str[0] == '.':
        if format_str.endswith('R'):
            if type(value) in six.integer_types:
                value = round_to_n(value, int(format_str[1]))
                format_str = ','
        if not format_str.endswith('G'):
            format_str = format_str + "G"
    try:
        value = format(value, format_str)
    except:
        value = format(value, '')

    return value



def process_position_parameter(param):
    """Process positioning parameters (left, top, width, height) given to
    df_to_table.

    If an integer, returns the right instance of the Cm class to allow it to be
    treated as cm. If missing, then default to 4cm. Otherwise, pass through
    whatever it gets.
    """
    if param is None:
        return Cm(4)
    elif type(param) is int:
        return Cm(param)
    else:
        return param



def df_to_table(slide, df, left=None, top=None, width=None, height=None,
                colnames=None, col_formatters=None, rounding=None,
                name=None, table_style=None):
    """Converts a Pandas DataFrame to a PowerPoint table on the given
    Slide of a PowerPoint presentation.

    The table is a standard Powerpoint table, and can easily be modified with
    the Powerpoint tools, for example: resizing columns, changing formatting etc.

    Parameters
    ----------
    slide: ``pptx.slide.Slide``
        slide object from the python-pptx library containing the slide on which
        you want the table to appear

    df: pandas ``DataFrame``
       DataFrame with the data

    left: int, optional
       Position of the left-side of the table, either as an integer in cm, or
       as an instance of a pptx.util Length class (pptx.util.Inches for
       example). Defaults to 4cm.

    top: int, optional
       Position of the top of the table, takes parameters as above.

    width: int, optional
       Width of the table, takes parameters as above.

    height: int, optional
       Height of the table, takes parameters as above.

    col_formatters: list, optional
       A n_columns element long list containing format specifications for each
       column. For example ['', ',', '.2'] does no special formatting for the
       first column, uses commas as thousands separators in the second column,
       and formats the third column as a float with 2 decimal places.

    rounding: list, optional
       A n_columns element long list containing a number for each integer
       column that requires rounding that is then multiplied by -1 and passed
       to round(). The practical upshot of this is that you can give something
       like ['', 3, ''], which does nothing for the 1st and 3rd columns (as
       they aren't integer values), but for the 2nd column, rounds away the 3
       right-hand digits (eg. taking 25437 to 25000).

    name: str, optional
       A name to be given to the table in the Powerpoint file. This is not
       displayed, but can help extract the table later to make further changes.
    
    table_style: str, optional
       Powerpoint table style to be used with generated table. These are strings
       that mirror the names of styles in Powerpoint, for example
            Medium Style 2 - Accent 1 -> table_style='MediumStyle2Accent1'

    Returns
    -------
    pptx.shapes.graphfrm.GraphicFrame
        The python-pptx table (GraphicFrame) object that was created (which can
        then be used to do further manipulation if desired)
    """
    left = process_position_parameter(left)
    top = process_position_parameter(top)
    width = process_position_parameter(width)
    height = process_position_parameter(height)

    rows, cols = df.shape
    shp = slide.shapes.add_table(rows+1, cols, left, top, width, height)

    if colnames is None:
        colnames = list(df.columns)

    # Insert the column names
    for col_index, col_name in enumerate(colnames):
        shp.table.cell(0,col_index).text = col_name

    m = df.values

    for row in range(rows):
        for col in range(cols):
            val = m[row, col]

            if col_formatters is None:
                text = str(val)
            else:
                text = _do_formatting(val, col_formatters[col])

            shp.table.cell(row+1, col).text = text

    if name is not None:
        shp.name = name

    if table_style is not None:
        table_graphic = shp._element.graphic.graphicData.tbl
        table_graphic[0][-1].text = getattr(TableStyle, table_style)

    return shp



def df_to_powerpoint(filename, df, **kwargs):
    """Converts a Pandas DataFrame to a table in a new, blank PowerPoint
    presentation.

    Creates a new PowerPoint presentation with the given filename, with a single
    slide containing a single table with the Pandas DataFrame data in it.

    The table is a standard Powerpoint table, and can easily be modified with
    the Powerpoint tools, for example: resizing columns, changing formatting
    etc.

    Parameters
    ----------
    filename: Filename to save the PowerPoint presentation as

    df: pandas ``DataFrame``
        DataFrame with the data

    **kwargs
        All other arguments that can be taken by ``df_to_table()`` (such as
        ``col_formatters`` or ``rounding``) can also be passed here.

    Returns
    -------
    pptx.shapes.graphfrm.GraphicFrame
        The python-pptx table (GraphicFrame) object that was created (which can
        then be used to do further manipulation if desired)
    """
    pres = Presentation()
    blank_slide_layout = pres.slide_layouts[6]
    slide = pres.slides.add_slide(blank_slide_layout)
    table = df_to_table(slide, df, **kwargs)
    pres.save(filename)

    return table
